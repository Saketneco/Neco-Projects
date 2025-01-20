from pptx import Presentation

# Create a presentation object
prs = Presentation()

# Function to add slide with title and content
def add_slide(title, content):
    slide_layout = prs.slide_layouts[1]  # Use 'Title and Content' layout
    slide = prs.slides.add_slide(slide_layout)
    
    title_placeholder = slide.shapes.title
    title_placeholder.text = title
    
    content_placeholder = slide.shapes.placeholders[1]
    content_placeholder.text = content

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]  # Title Slide Layout
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Overview of TAMUQ Information Technology Department"
subtitle.text = "Administrative, Business, and IT Resources\nYour Name\nDate"

# Slide 2: TAMUQ Information Technology Department Overview
add_slide(
    "TAMUQ Information Technology Department Overview",
    "• Introduction to TAMUQ IT Department.\n• Mission and vision.\n• Supporting students, faculty, and staff.\n• Aligning with TAMU global IT standards."
)

# Slide 3: Administrative and Business IT Support
add_slide(
    "Administrative and Business IT Support",
    "• IT support for administrative functions (HR, finance).\n• Integration with TAMU's global admin network.\n• Efficiency, security, and scalability of systems."
)

# Slide 4: Communication and Collaboration
add_slide(
    "Communication and Collaboration",
    "• Platforms for communication: email, chat, video conferencing.\n• Collaborative tools: Teams, Zoom, SharePoint.\n• Supporting remote work and learning."
)

# Slide 5: End-Point Computing
add_slide(
    "End-Point Computing",
    "• Devices supported: laptops, desktops, tablets.\n• Ensuring secure access for all users.\n• Support for software installation and troubleshooting."
)

# Slide 6: Infrastructure
add_slide(
    "Infrastructure",
    "• IT infrastructure: data centers, networks, servers.\n• High availability and redundancy strategies.\n• Cloud computing and sustainable infrastructure."
)

# Slide 7: Security
add_slide(
    "Security",
    "• Cybersecurity measures for data, networks, and devices.\n• Security policies for staff and students.\n• Ongoing security training and awareness."
)

# Slide 8: Teaching and Learning IT Support
add_slide(
    "Teaching and Learning IT Support",
    "• IT resources for teaching and learning: classroom technologies.\n• E-learning tools: Blackboard, Moodle.\n• IT support for research and academic initiatives."
)

# Slide 9: TAMUQ IT Professional Services
add_slide(
    "TAMUQ IT Professional Services",
    "• IT professional services for faculty and staff.\n• Support for teaching, research, and administration.\n• Custom IT solutions and consulting."
)

# Slide 10: TAMU Resources for IT
add_slide(
    "TAMU Resources for IT",
    "• Overview of IT resources at TAMU.\n• Access to TAMU’s IT systems and tools.\n• Collaboration between TAMU and TAMUQ for enhanced IT solutions."
)

# Slide 11: TAMUQ IT Resources
add_slide(
    "TAMUQ IT Resources",
    "• Local IT resources and support at TAMUQ.\n• IT help desk and technical support.\n• Access to databases, software, and research tools."
)

# Slide 12: Conclusion
add_slide(
    "Conclusion",
    "• Summary of TAMUQ IT services and resources.\n• IT's role in supporting administrative and academic functions.\n• Commitment to innovation and user support."
)

# Slide 13: Q&A
add_slide(
    "Q&A",
    "• Invite questions from the audience.\n• Open discussion on TAMUQ IT services."
)

# Save presentation
prs.save("D:\\USER PROFILE DATA\\Desktop\\Project-2(Mail Automation)\\Data\\Department_Presentation.pptx")
